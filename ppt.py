generate_calculator_ppt:any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create Presentation
prs = Presentation()

# Function to add a slide
def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    text_box = slide.placeholders[1]
    text_frame = text_box.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = content
    p.font.size = Pt(18)
    p.font.name = 'Calibri'
    p.font.color.rgb = RGBColor(0, 0, 0)

# Slide contents
slides = [
    ("Advanced Calculator in Python", "A modern Python project developed using Visual Studio Code\nPresented by: Your Name"),
    ("Project Overview", "A modern and efficient Advanced Calculator using Python.\nHandles both basic and scientific operations."),
    ("Objectives", "• Create a user-friendly calculator\n• Support basic and advanced operations\n• Use Python math library\n• Console-based program"),
    ("Features", "Basic Operations: Addition, Subtraction, Multiplication, Division\nAdvanced Operations: Power, Square Root, sin, cos, tan, log, log10\nError handling included"),
    ("Tools Used", "• Python 3.8+\n• Visual Studio Code\n• Python Math Library\n• Command Line Interface"),
    ("Workflow", "User Input → Menu Selection → Calculation → Display Result"),
    ("Python Code Sample", "import math\n\ndef calculator():\n    print('===== ADVANCED CALCULATOR =====')\n    # ... rest of the code ..."),
    ("Sample Output", "===== ADVANCED CALCULATOR =====\n1. Addition\n2. Subtraction\nEnter your choice: 1\nEnter first number: 10\nEnter second number: 5\nResult: 15"),
    ("Applications", "• Educational tool\n• Scientific calculations\n• Quick Python project\n• School/College submission"),
    ("Conclusion", "Python के साथ advanced calculator बनाना सरल और प्रभावी है।\nProject logical thinking और programming skills को बढ़ाता है।"),
    ("Thank You", "Questions?")
]

# Add slides
for title, content in slides:
    add_slide(title, content)

# Save presentation
ppt_file = "Advanced_Calculator_Python.pptx"
prs.save(ppt_file)
print(f"PPT successfully created: {ppt_file}")